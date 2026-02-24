# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx",
# ]
# ///
"""Extract slide content and notes from a PPTX file."""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def extract_slide_data(pptx_path: str) -> list[dict]:
    """Extract text content and notes from each slide."""
    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': i,
            'title': '',
            'texts': [],
            'notes': '',
            'has_images': False,
            'image_count': 0,
            'shapes': [],
        }

        # Extract text from shapes
        for shape in slide.shapes:
            shape_info = {
                'type': str(shape.shape_type),
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
            }

            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Check for bullet level
                        level = para.level if para.level else 0
                        bold = any(run.font.bold for run in para.runs if run.font.bold)
                        paragraphs.append({
                            'text': text,
                            'level': level,
                            'bold': bold,
                        })
                if paragraphs:
                    shape_info['paragraphs'] = paragraphs

                    # First text frame with content might be the title
                    if not slide_info['title'] and shape.shape_type == 13:  # Title
                        slide_info['title'] = paragraphs[0]['text']

            if hasattr(shape, 'image'):
                slide_info['has_images'] = True
                slide_info['image_count'] += 1
                shape_info['has_image'] = True
                shape_info['image_content_type'] = shape.image.content_type

            slide_info['shapes'].append(shape_info)

        # Try to get title from placeholder
        if not slide_info['title']:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 0:  # Title placeholder
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        slide_info['title'] = shape.text_frame.text.strip()
                        break

        # Collect all text content
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_info['texts'].append(text)

        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_info['notes'] = notes_slide.notes_text_frame.text.strip()

        slides_data.append(slide_info)

    return slides_data


def main() -> None:
    """Run the extraction."""
    if len(sys.argv) < 2:
        print('Usage: uv run extract_pptx.py <pptx_file> [output_json]')
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    data = extract_slide_data(pptx_path)

    output = json.dumps(data, ensure_ascii=False, indent=2)

    if output_path:
        Path(output_path).write_text(output, encoding='utf-8')
        print(f'Extracted {len(data)} slides to {output_path}')
    else:
        print(output)


if __name__ == '__main__':
    main()
