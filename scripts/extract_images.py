# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx",
# ]
# ///
"""Extract images from a PPTX file, organized by slide number."""

import sys
from pathlib import Path

from pptx import Presentation


def extract_images(pptx_path: str, output_dir: str) -> None:
    """Extract all images from PPTX, saving per slide."""
    prs = Presentation(pptx_path)
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)

    for i, slide in enumerate(prs.slides, 1):
        img_idx = 0
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                img_idx += 1
                ext = shape.image.content_type.split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'
                filename = f'slide{i:02d}_{img_idx}.{ext}'
                filepath = out / filename
                filepath.write_bytes(shape.image.blob)
                print(f'Slide {i}: saved {filepath}')

    if img_idx == 0:
        print('No images found in the last slide checked.')


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print('Usage: uv run extract_images.py <pptx_file> <output_dir>')
        sys.exit(1)

    extract_images(sys.argv[1], sys.argv[2])
